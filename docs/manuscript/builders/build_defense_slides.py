"""
build_defense_slides.py  —  SmartSpend Pre-Final Defense Presentation
16 slides, 8-9 minute flow, Lorma maroon/gold theme using python-pptx.
Run:  python build_defense_slides.py
Out:  docs/manuscript/output/SmartSpend_Defense_Slides.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / '..' / 'output' / 'SmartSpend_Defense_Slides.pptx'

# ── Colours ────────────────────────────────────────────────────────────────────
MAROON  = RGBColor(0x5C, 0x0E, 0x24)
GOLD    = RGBColor(0xC9, 0xA8, 0x4C)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x22, 0x22, 0x22)
LGREY   = RGBColor(0xF0, 0xF0, 0xF0)
DGREY   = RGBColor(0x55, 0x55, 0x55)
BLUE    = RGBColor(0x2A, 0x4A, 0x7F)
GREEN   = RGBColor(0x1A, 0x6B, 0x3A)
AMBER   = RGBColor(0xE6, 0x5C, 0x00)

FONT    = 'Calibri'

# ── Slide dimensions: Widescreen 16:9 ─────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from lxml import etree


def blank_layout():
    return prs.slide_layouts[6]  # blank


def add_slide():
    return prs.slides.add_slide(blank_layout())


def rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def txt(slide, text, left, top, width, height,
        size=24, bold=False, color=WHITE, align='left',
        italic=False, wrap=True, font=None):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = {'left':   PP_ALIGN.LEFT,
                   'center': PP_ALIGN.CENTER,
                   'right':  PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = str(text)
    run.font.name   = font or FONT
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def bullet_box(slide, items, left, top, width, height,
               size=18, color=DARK, marker='●', gap=0.08,
               bold_first=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(int(gap * 72))
        run = p.add_run()
        run.text = f'{marker}  {item}'
        run.font.name  = FONT
        run.font.size  = Pt(size)
        run.font.color.rgb = color
        run.font.bold  = (bold_first and i == 0)
    return txb


def accent_bar(slide, top=0.0, height=0.08, color=GOLD):
    """Full-width accent bar."""
    rect(slide, 0, top, 13.33, height, color)


def card(slide, left, top, width, height, fill='F5F5F5', border_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    if border_color:
        shape.line.color.rgb = RGBColor.from_string(border_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def slide_number(slide, n, total=16):
    txt(slide, f'{n} / {total}', 12.5, 7.1, 0.8, 0.35,
        size=11, color=DGREY, align='right')


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()

# Full background maroon
rect(sl, 0, 0, 13.33, 7.5, MAROON)
# Gold accent bar bottom
rect(sl, 0, 6.9, 13.33, 0.6, GOLD)

txt(sl, 'SmartSpend', 0.8, 0.9, 11.5, 1.5,
    size=60, bold=True, color=WHITE, align='center')
txt(sl, 'An AI-Assisted Mobile Financial Tracking and Advisory Application\nfor Personal Financial Management',
    0.8, 2.3, 11.5, 1.2, size=22, color=GOLD, align='center', italic=True)

# Horizontal rule
rect(sl, 2.5, 3.65, 8.3, 0.04, GOLD)

txt(sl, 'Pre-Final Defense Presentation', 0.8, 3.8, 11.5, 0.5,
    size=17, color=WHITE, align='center', bold=True)
txt(sl, 'Brix A. Directo  |  Cyrille John M. Rubis  |  Djaunathan Albert S. Madayag',
    0.8, 4.35, 11.5, 0.45, size=14, color=LGREY, align='center')
txt(sl, 'Lucid Frame  |  BSIT 4th Year  |  Lorma Colleges CCSE  |  AY 2026–2027, 1st Semester',
    0.8, 4.82, 11.5, 0.4, size=12, color=LGREY, align='center')
txt(sl, 'Adviser: Janelli M. Mendez, DIT', 0.8, 5.25, 11.5, 0.35,
    size=12, color=GOLD, align='center')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'The Problem', 0.4, 0.12, 12.5, 0.75, size=32, bold=True, color=WHITE)
slide_number(sl, 2)

# 3 stat cards
stats = [
    ('50%', 'Filipino adults with\nbank accounts\n(BSP CFIS, 2025)', 'EAD9E0'),
    ('74%', 'Filipinos who answered\nfinancial literacy questions\ncorrectly (BSP, 2025)', 'DAE8F5'),
    ('18%', 'Global consumers using AI\nspecifically for budgeting\n(EY Survey, 2026)', 'D9F0DA'),
]
for i, (stat, label, color) in enumerate(stats):
    x = 0.6 + i * 4.1
    card(sl, x, 1.3, 3.7, 2.4, fill=color, border_color='CCCCCC')
    txt(sl, stat, x, 1.35, 3.7, 1.1, size=48, bold=True, color=MAROON, align='center')
    txt(sl, label, x, 2.4, 3.7, 1.2, size=14, color=DARK, align='center')

txt(sl, 'The Gap:', 0.5, 3.9, 12.0, 0.45, size=20, bold=True, color=MAROON)
txt(sl,
    'Most Filipinos do not track finances — not because they don\'t want to, but because '
    'traditional tools (spreadsheets, manual apps) are too tedious. '
    'No existing free Android app combines agentic AI + offline capability + Filipino-English + a Financial Health Score.',
    0.5, 4.35, 12.3, 1.0, size=17, color=DARK)

card(sl, 0.5, 5.45, 12.3, 0.9, fill='F5EEF0', border_color='CC9999')
txt(sl, '📌  SmartSpend removes friction: say "I spent 85 pesos on lunch" — it logs automatically. No forms. No dropdowns.',
    0.65, 5.52, 12.0, 0.75, size=16, bold=True, color=MAROON)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RESEARCH OBJECTIVES
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'Research Objectives', 0.4, 0.12, 12.5, 0.75, size=32, bold=True, color=WHITE)
slide_number(sl, 3)

obj_data = [
    ('1', 'Assessment',
     'Assess the existing financial management practices, common budgeting challenges, and '
     'expense tracking behaviors of parents (35–55) and young professionals (21–35) in La Union.',
     'EAD9E0', MAROON),
    ('2', 'Development',
     'Design and develop the SmartSpend mobile application, including a comparative evaluation '
     'of 15 LLM API providers to select the optimal model for Filipino-English agentic AI.',
     'D9E8F5', BLUE),
    ('3', 'Evaluation',
     'Evaluate the usability of SmartSpend using the System Usability Scale (SUS) with '
     '30 purposively selected respondents. Target: SUS Score ≥ 80 (Good).',
     'D9F0DA', GREEN),
]
for i, (num, label, desc, bg, col) in enumerate(obj_data):
    y = 1.25 + i * 1.8
    card(sl, 0.5, y, 12.3, 1.6, fill=bg, border_color='CCCCCC')
    card(sl, 0.5, y, 0.75, 1.6, fill=f'{col[0]:02X}{col[1]:02X}{col[2]:02X}')
    txt(sl, num, 0.5, y + 0.3, 0.75, 1.0, size=36, bold=True, color=WHITE, align='center')
    txt(sl, f'Objective {num} — {label}', 1.35, y + 0.1, 11.0, 0.55,
        size=17, bold=True, color=col)
    txt(sl, desc, 1.35, y + 0.65, 11.0, 0.85, size=14, color=DARK)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT IS SMARTSPEND?
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'What is SmartSpend?', 0.4, 0.12, 12.5, 0.75, size=32, bold=True, color=WHITE)
slide_number(sl, 4)

txt(sl, 'User Input  →  AI Parsing  →  Autonomous Action  →  SQLite  →  Analytics & Insights',
    0.5, 1.1, 12.3, 0.5, size=16, bold=True, color=MAROON, align='center')

features = [
    ('🤖', '31 Agentic AI\nActions', 'AI takes real actions\non your financial data'),
    ('📊', 'Financial Health\nScore 0–100', 'Behavioral metric,\ndual-mode formula'),
    ('🎤', 'Multi-Modal Input', 'Voice, OCR, Barcode,\nBatch Screenshots (40+)'),
    ('📴', 'Offline-First', 'SQLite + optional\nFirebase sync'),
    ('🇵🇭', 'Filipino-First', 'Taglish AI, GCash,\nSSS, PhilHealth, BIR'),
    ('🆓', 'Always Free', 'Gemini + Groq +\nCerebras free tier'),
]
for i, (emoji, title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 1.8 + row * 2.4
    card(sl, x, y, 3.9, 2.1, fill='FAFAFA', border_color='DDDDDD')
    txt(sl, emoji, x, y + 0.1, 3.9, 0.75, size=28, align='center')
    txt(sl, title, x, y + 0.75, 3.9, 0.65, size=15, bold=True, color=MAROON, align='center')
    txt(sl, desc, x, y + 1.35, 3.9, 0.65, size=12, color=DGREY, align='center')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECHNOLOGY STACK
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'Technology Stack', 0.4, 0.12, 12.5, 0.75, size=32, bold=True, color=WHITE)
slide_number(sl, 5)

stack = [
    ('Flutter / Dart', 'App Framework', 'Android (arm64-v8a), 45 MB, single codebase', 'E8F0FE'),
    ('SQLite v11', 'Local Database', '20 tables, offline-first, full CRUD without internet', 'FFF8E1'),
    ('Firebase', 'Cloud Layer', 'Firestore sync + Auth (Google/email) + Crashlytics + Remote Config', 'E8F5E9'),
    ('Gemini 3.1 Flash-Lite', 'Primary LLM', '1,000 req/day free — best Filipino-English, 1M token context', 'F3E5F5'),
    ('Groq LLaMA 3.3 70B', 'Fallback LLM 2', '14,400 req/day free — best open-source reasoning', 'E3F2FD'),
    ('Google ML Kit', 'OCR + Barcode', 'On-device text recognition, no API key, works offline', 'FBE9E7'),
]
for i, (tech, role, detail, color) in enumerate(stack):
    y = 1.15 + i * 1.02
    card(sl, 0.4, y, 12.5, 0.94, fill=color, border_color='DDDDDD')
    txt(sl, tech, 0.55, y + 0.1, 2.4, 0.72, size=15, bold=True, color=MAROON)
    txt(sl, role, 3.0, y + 0.1, 2.0, 0.35, size=11, bold=True, color=BLUE)
    txt(sl, detail, 3.0, y + 0.45, 9.7, 0.42, size=12, color=DARK)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — AGENTIC AI ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'Agentic AI Architecture', 0.4, 0.12, 12.5, 0.75, size=32, bold=True, color=WHITE)
slide_number(sl, 6)

# Flow diagram
flow_items = [
    ('User Input\n(Voice / Text / Camera / Screenshot)', '2A4A7F'),
    ('AI Parsing\n(Gemini / Groq / Cerebras)', '5C0E24'),
    ('31 Action Types\n(Intent Classification)', 'C97000'),
    ('SQLite Write\n(Autonomous Action)', '1A6B3A'),
    ('FHS Update\n(score_service.dart)', '8B1A6B'),
]
for i, (label, color) in enumerate(flow_items):
    x = 0.35 + i * 2.58
    card(sl, x, 1.2, 2.3, 1.35, fill=color, border_color='FFFFFF')
    txt(sl, label, x, 1.2, 2.3, 1.35, size=13, bold=True, color=WHITE, align='center')
    if i < len(flow_items) - 1:
        txt(sl, '→', x + 2.3, 1.6, 0.28, 0.55, size=24, bold=True, color=MAROON, align='center')

txt(sl, 'Perceive  →  Decide  →  Act  (genuine agentic loop)',
    0.5, 2.75, 12.3, 0.45, size=16, bold=True, color=MAROON, align='center')

# 5-provider chain
txt(sl, '5-Provider Auto-Failover (all free tier):', 0.5, 3.3, 12.0, 0.4,
    size=15, bold=True, color=BLUE)
providers = [
    ('1. Gemini 3.1\nFlash-Lite', '1,000/day\n(primary)', 'F3E5F5'),
    ('2. Gemini 3.5\nFlash', '250/day\n(fallback 1)', 'EDE7F6'),
    ('3. Groq LLaMA\n3.3 70B', '14,400/day\n(fallback 2)', 'E3F2FD'),
    ('4. Groq LLaMA\n3.1 8B', '14,400/day\n(fallback 3)', 'E8F5E9'),
    ('5. Cerebras\nLLaMA 3.1', '1M tokens/day\n(fallback 4)', 'FFF8E1'),
]
for i, (name, quota, color) in enumerate(providers):
    x = 0.5 + i * 2.55
    card(sl, x, 3.8, 2.3, 1.35, fill=color, border_color='CCCCCC')
    txt(sl, name, x, 3.82, 2.3, 0.72, size=13, bold=True, color=MAROON, align='center')
    txt(sl, quota, x, 4.55, 2.3, 0.55, size=11, color=BLUE, align='center')

card(sl, 0.5, 5.3, 12.3, 0.9, fill='FFF3F5', border_color='CC9999')
txt(sl, '💡  Why Context Injection (not RAG)?  Per-user data (~5K tokens) fits in one prompt. '
    'RAG adds vector-search overhead designed for thousands of documents — unnecessary here.',
    0.65, 5.37, 12.0, 0.75, size=14, color=MAROON)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — LLM BENCHMARKING (COMPARATIVE STUDY)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'LLM Comparative Benchmarking', 0.4, 0.12, 12.5, 0.75, size=32, bold=True, color=WHITE)
slide_number(sl, 7)

txt(sl, 'Evaluation of 15 LLM APIs — Selection Criteria Weighted for Filipino-English Mobile Deployment',
    0.5, 1.05, 12.3, 0.4, size=15, italic=True, color=DGREY)

# Criteria pills
criteria = [
    ('Filipino-English Accuracy', '25%', 'EAD9E0'),
    ('Speed / Latency',            '20%', 'DAE8F5'),
    ('Tool Use / JSON',            '20%', 'D9F0DA'),
    ('Free Tier',                  '15%', 'FFF8DC'),
    ('Context Window',             '10%', 'F0E6FF'),
    ('Financial Reasoning',        '10%', 'FFE8D0'),
]
for i, (label, weight, color) in enumerate(criteria):
    x = 0.4 + (i % 3) * 4.3
    y = 1.55 + (i // 3) * 0.72
    card(sl, x, y, 4.0, 0.62, fill=color, border_color='CCCCCC')
    txt(sl, f'{weight}  {label}', x + 0.1, y + 0.1, 3.8, 0.42,
        size=14, bold=True, color=MAROON)

# Top 5 selected models table
txt(sl, 'Selected Models (Top 5):', 0.5, 3.1, 12.0, 0.38, size=15, bold=True, color=BLUE)

models_data = [
    ('Model',               'Provider',  'Speed (t/s)', 'Filipino', 'Tool Use', 'Free Tier',    'Role'),
    ('Gemini 3.1 Flash-Lite','Google',   '~400–600',    '★★★★★',    '★★★★★',   '1,000/day',    '✅ PRIMARY'),
    ('Gemini 3.5 Flash',    'Google',    '~200–400',    '★★★★★',    '★★★★★',   '250/day',      '✅ Fallback 1'),
    ('LLaMA 3.3 70B',       'Groq LPU',  '~315',        '★★★★☆',    '★★★★★',   '14,400/day',   '✅ Fallback 2'),
    ('LLaMA 3.1 8B',        'Groq LPU',  '~800',        '★★★★☆',    '★★★★☆',   '14,400/day',   '✅ Fallback 3'),
    ('LLaMA 3.1 70B',       'Cerebras',  '~1,800',      '★★★★☆',    '★★★★☆',   '1M tokens/day','✅ Fallback 4'),
]
COL_W = [2.5, 1.2, 1.0, 0.9, 0.9, 1.3, 1.4]
COL_X = [0.4]
for w in COL_W[:-1]: COL_X.append(COL_X[-1] + w)

for ri, row in enumerate(models_data):
    y = 3.55 + ri * 0.58
    for ci, (cell_text, cw, cx) in enumerate(zip(row, COL_W, COL_X)):
        bg = '5C0E24' if ri == 0 else ('FFF3E0' if ri == 1 else 'FAFAFA')
        card(sl, cx, y, cw - 0.04, 0.52, fill=bg,
             border_color='AAAAAA' if ri > 0 else None)
        fc = WHITE if ri == 0 else (MAROON if ri == 1 else DARK)
        txt(sl, cell_text, cx + 0.05, y + 0.05, cw - 0.1, 0.42,
            size=11 if ri > 0 else 10, bold=(ri == 0 or ri == 1), color=fc)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GENERAL vs FINANCE-SPECIALIZED LLMs
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'General LLMs vs Finance-Specialized LLMs', 0.4, 0.12, 12.5, 0.75,
    size=28, bold=True, color=WHITE)
slide_number(sl, 8)

txt(sl, 'Why FinGPT, Fin-R1, and BloombergGPT are NOT the right choice for SmartSpend',
    0.5, 1.05, 12.3, 0.38, size=14, italic=True, color=DGREY)

comparisons = [
    ('Task', 'Finance-Specialized LLMs\n(FinGPT, Fin-R1, BloombergGPT)', 'General Multilingual LLMs\n(Gemini 3.1, LLaMA 3.3 70B)'),
    ('Filipino-English\nExpense Parsing',
     '❌  Not trained on Tagalog/Filipino.\nPoor on local merchant names.',
     '✅  Google training covers SEA languages.\nHandles Taglish correctly.'),
    ('Agentic JSON\nOutput (31 actions)',
     '❌  No function calling support.\nBloombergGPT has no API.',
     '✅  Native function calling.\nReliable JSON output.'),
    ('PH Financial\nAdvisory (SSS, BIR)',
     '❌  Trained on US market data.\nNo PH government knowledge.',
     '✅  Gemini includes PH gov sources.\nCorrect SSS/PhilHealth rates.'),
    ('Free API Access',
     '❌  Self-hosted only (needs GPU).\nNo hosted free API.',
     '✅  Gemini: 1,000/day free.\nGroq: 14,400/day free.'),
]
col_w = [1.9, 5.2, 5.1]
col_x = [0.4, 2.35, 7.6]
row_shades = ['5C0E24', 'F5EEF0', 'FAFAFA', 'F5EEF0', 'FAFAFA']
for ri, (row, rshade) in enumerate(zip(comparisons, row_shades)):
    y = 1.52 + ri * 1.1
    for ci, (cell_text, cw, cx) in enumerate(zip(row, col_w, col_x)):
        bg = rshade if ri > 0 else '5C0E24'
        if ri > 0 and ci == 1: bg = 'FFEBEE'
        if ri > 0 and ci == 2: bg = 'E8F5E9'
        card(sl, cx, y, cw - 0.05, 1.0, fill=bg, border_color='BBBBBB')
        fc = WHITE if ri == 0 else (RGBColor(0x8B,0x00,0x00) if ci == 1 and ri > 0 else
                                    (RGBColor(0x00,0x60,0x00) if ci == 2 and ri > 0 else DARK))
        txt(sl, cell_text, cx + 0.08, y + 0.08, cw - 0.16, 0.84,
            size=12 if ri > 0 else 11, bold=(ri == 0), color=fc)

card(sl, 0.4, 7.05, 12.5, 0.35, fill='FFF3F5', border_color='CC9999')
txt(sl, '📌  Conclusion (FrontierFinance Benchmark, 2026): Tool harness architecture '
    'matters more than model specialization. General multilingual LLMs outperform '
    'finance-specialized models on SmartSpend\'s Filipino-English consumer workload.',
    0.55, 7.08, 12.2, 0.28, size=11, color=MAROON, bold=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FINANCIAL HEALTH SCORE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'Financial Health Score (FHS)', 0.4, 0.12, 12.5, 0.75, size=32, bold=True, color=WHITE)
slide_number(sl, 9)

txt(sl, '0 to 100 behavioral metric — computed from user transaction data — no surveys, no bank connectivity',
    0.5, 1.08, 12.3, 0.38, size=15, italic=True, color=DGREY)

# Two mode columns
rect(sl, 0.4, 1.55, 5.8, 0.45, BLUE)
txt(sl, 'Full Mode  (Income Tracking ON)', 0.4, 1.55, 5.8, 0.45,
    size=14, bold=True, color=WHITE, align='center')
rect(sl, 6.5, 1.55, 6.4, 0.45, GREEN)
txt(sl, 'Lightweight Mode  (Income Tracking OFF)', 6.5, 1.55, 6.4, 0.45,
    size=14, bold=True, color=WHITE, align='center')

full_comps = [
    ('Savings Rate', '25 pts', '25 × min(1, savingsRate / 0.20)', 'Warren & Tyagi (2005) 50/30/20 rule'),
    ('Overspend Control', '25 pts', '25 × (1 − overDays / activeDays)', 'FinHealth Score® Spend pillar'),
    ('Budget Adherence', '25 pts', '25 × (onBudget / totalBudgets)', 'Ramsey (2003) zero-based budgeting'),
    ('Logging Consistency', '25 pts', '25 × (loggedDays / activeDays)', 'Thaler & Sunstein (2008) nudge theory'),
]
lite_comps = [
    ('Spending Restraint', '25 pts', 'vs user-set spending limit'),
    ('Logging Consistency', '25 pts', 'same formula as Full Mode'),
    ('Category Balance', '25 pts', 'no single category > 40%'),
    ('Habit Streak', '25 pts', '14-day streak = full score'),
]

for i, (name, pts, formula, basis) in enumerate(full_comps):
    y = 2.1 + i * 1.1
    card(sl, 0.4, y, 5.8, 1.0, fill='EEF4FF', border_color='AABBDD')
    txt(sl, f'{name}  ({pts})', 0.55, y + 0.05, 5.5, 0.4, size=14, bold=True, color=BLUE)
    txt(sl, formula, 0.55, y + 0.43, 5.5, 0.3, size=11, color=DARK, font='Courier New')
    txt(sl, f'Basis: {basis}', 0.55, y + 0.72, 5.5, 0.25, size=10, italic=True, color=DGREY)

for i, (name, pts, desc) in enumerate(lite_comps):
    y = 2.1 + i * 1.1
    card(sl, 6.5, y, 6.4, 1.0, fill='EFFFEF', border_color='AADDAA')
    txt(sl, f'{name}  ({pts})', 6.65, y + 0.05, 6.1, 0.4, size=14, bold=True, color=GREEN)
    txt(sl, desc, 6.65, y + 0.52, 6.1, 0.38, size=13, color=DARK)

# Adjustments bar
card(sl, 0.4, 6.55, 12.5, 0.8, fill='FFF0F3', border_color='CC9999')
txt(sl, '⚠️  Score Adjustments:  '
    'Warning Decay (−5 pts/day, max −15) — ignoring budget warnings  |  '
    'Gap Adjustment (−3 or +2 pts/day) — unlogged days confirmed by user',
    0.55, 6.62, 12.2, 0.66, size=13, color=MAROON, bold=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — APP FEATURE COMPARISON
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'Comparative Study — SmartSpend vs Existing Apps', 0.4, 0.12, 12.5, 0.75,
    size=28, bold=True, color=WHITE)
slide_number(sl, 10)

txt(sl, '22 apps reviewed (14 international, 8 Philippine-context)  |  August 2026',
    0.5, 1.05, 12.3, 0.35, size=14, italic=True, color=DGREY)

headers = ['Feature', 'SmartSpend', 'YNAB', 'Monarch', 'BudgetPH', 'Alkansya AI', 'Pera Coach']
rows_data = [
    ['31 Agentic AI Actions',   '✅ 31',     '❌',    '❌',    '❌ Insights','❌',       '❌'],
    ['Filipino-English AI',     '✅ Full',   '❌',    '❌',    '❌',         '✅ iOS',   '✅'],
    ['Financial Health Score',  '✅ 0–100',  '❌',    '❌',    '✅ Simpler', '✅ Web',   '❌'],
    ['Offline Mode',            '✅ Full',   '❌',    '❌',    '✅',         '❌',       '❌'],
    ['Batch Screenshots (40+)', '✅',        '❌',    '❌',    '❌',         '❌',       '❌'],
    ['Voice Input',             '✅',        '❌',    '❌',    '❌',         '❌',       '❌'],
    ['Free (Android)',          '✅ Always', '❌',    '❌',    '✅ PWA',     '❌ iOS+$', '✅ GCash'],
    ['Gamification (Badges)',   '✅ 23',     '❌',    '❌',    '✅ XP/lvl',  '❌',       '❌'],
]
col_w2 = [2.3, 1.5, 1.3, 1.5, 1.5, 1.6, 1.5]
col_x2 = [0.4]; [col_x2.append(col_x2[-1] + w) for w in col_w2[:-1]]

for ri, row in enumerate([headers] + rows_data):
    y = 1.5 + ri * 0.63
    for ci, (cell_text, cw, cx) in enumerate(zip(row, col_w2, col_x2)):
        is_hdr = ri == 0
        is_ss  = ci == 1
        bg = ('5C0E24' if is_hdr else
              'FFF0E8' if is_ss else
              'F9F9F9' if ri % 2 == 0 else 'FFFFFF')
        card(sl, cx, y, cw - 0.04, 0.57, fill=bg, border_color='BBBBBB')
        fc = (WHITE if is_hdr else
              MAROON if is_ss else
              (GREEN if '✅' in str(cell_text) else
               (RGBColor(0xBB,0x00,0x00) if '❌' in str(cell_text) else DARK)))
        txt(sl, cell_text, cx + 0.06, y + 0.08, cw - 0.12, 0.4,
            size=11 if ri > 0 else 10, bold=(ri == 0 or (ri > 0 and ci == 1)), color=fc)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SMART IMPORT & MULTI-MODAL INPUT
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'Multi-Modal Input & Smart Import', 0.4, 0.12, 12.5, 0.75, size=32, bold=True, color=WHITE)
slide_number(sl, 11)

modes = [
    ('📹\nLive Camera', 'Barcode/QR live detection\n+ Receipt mode OCR\n+ Torch toggle', '2A4A7F'),
    ('🖼️\nSingle Photo', 'Auto-routes: barcode,\nreceipt, or screenshot\nbased on content', '5C0E24'),
    ('📱\nBatch Screenshots', 'Up to 10 images\nAuto-detects 40+ platforms\nAI prompt per platform', '1A6B3A'),
    ('📋\nPaste Text', 'GCash/BPI/Maya\nbank export text\nOCR text import', 'C97000'),
]
for i, (title, desc, color) in enumerate(modes):
    x = 0.4 + i * 3.2
    card(sl, x, 1.2, 3.0, 2.5, fill=color)
    txt(sl, title, x, 1.25, 3.0, 1.0, size=18, bold=True, color=WHITE, align='center')
    txt(sl, desc, x, 2.22, 3.0, 1.4, size=13, color=WHITE, align='center')

txt(sl, '40+ platforms auto-detected with platform-specific AI extraction prompts:',
    0.5, 3.9, 12.3, 0.38, size=15, bold=True, color=MAROON)

platforms = [
    '🛍️  Shopee, Lazada, TikTok Shop, Carousell, Amazon, AliExpress',
    '🎮  Steam, Google Play, App Store, Codashop, Garena, MLBB, Genshin, Valorant',
    '💳  GCash, Maya, GrabPay, ShopeePay, Coins.ph, PayPal, BPI, BDO, UnionBank',
    '🚗  GrabFood, Foodpanda, Grab rides, Angkas, Maxim, Lalamove',
    '🎬  Netflix, Spotify, YouTube Premium, Disney+, Viu, Vivamax',
]
for i, p in enumerate(platforms):
    y = 4.35 + i * 0.56
    card(sl, 0.5, y, 12.3, 0.5, fill='F8F8F8', border_color='DDDDDD')
    txt(sl, p, 0.65, y + 0.07, 12.0, 0.36, size=13, color=DARK)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — METHODOLOGY
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'Research Methodology', 0.4, 0.12, 12.5, 0.75, size=32, bold=True, color=WHITE)
slide_number(sl, 12)

# Kanban phases
phases = [
    ('Backlog', '1', 'EAD9E0'), ('Requirements', '2', 'DAE8F5'),
    ('Design', '3', 'D9EAF5'), ('Development', '4', 'D9F0DA'),
    ('Testing', '5', 'FFF8DC'), ('Deployment', '6', 'FFE8D0'),
    ('Done/Review', '7', 'F0E6FF'),
]
for i, (phase, num, color) in enumerate(phases):
    x = 0.35 + i * 1.84
    card(sl, x, 1.2, 1.74, 0.55, fill='5C0E24')
    txt(sl, f'{num}. {phase}', x, 1.2, 1.74, 0.55,
        size=10, bold=True, color=WHITE, align='center')
    card(sl, x, 1.78, 1.74, 0.55, fill=color, border_color='CCCCCC')

txt(sl, 'Agile Kanban Methodology — 7 phases, each producing a deliverable',
    0.5, 2.45, 12.3, 0.35, size=13, italic=True, color=DGREY, align='center')

# Research design cards
design_items = [
    ('Research Design', 'Mixed Methods\n(Developmental + Descriptive)', 'E8F0FE'),
    ('Population', '30 Respondents:\n20 Parents (35–55)\n10 Young Professionals (21–35)', 'FFF0F3'),
    ('Sampling', 'Purposive Sampling\nLa Union, Philippines', 'E8F5E9'),
    ('Data Tools', 'Structured Survey\n(Obj. 1)  +  SUS (Obj. 3)\n+  Expert Validators (Obj. 2)', 'FFF8DC'),
    ('Evaluation', 'System Usability Scale\n(Brooke, 1996)\nTarget: SUS ≥ 80', 'F3E5F5'),
    ('Framework', 'IPO Model\n(Input–Process–Output)\nConceptual Framework', 'E0F7FA'),
]
for i, (title, content, color) in enumerate(design_items):
    col = i % 3; row = i // 3
    x = 0.4 + col * 4.25
    y = 3.05 + row * 2.15
    card(sl, x, y, 4.05, 2.0, fill=color, border_color='CCCCCC')
    txt(sl, title, x + 0.1, y + 0.1, 3.85, 0.45, size=13, bold=True, color=MAROON)
    txt(sl, content, x + 0.1, y + 0.58, 3.85, 1.3, size=13, color=DARK)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SUS EVALUATION (placeholder for results)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'Usability Evaluation — SUS Results', 0.4, 0.12, 12.5, 0.75,
    size=32, bold=True, color=WHITE)
slide_number(sl, 13)

# Result placeholder card
card(sl, 3.5, 1.2, 6.3, 2.8, fill='F5EEF0', border_color='5C0E24')
txt(sl, 'Overall SUS Score', 3.5, 1.3, 6.3, 0.55, size=18, bold=True,
    color=MAROON, align='center')
txt(sl, '[ _____ / 100 ]', 3.5, 1.9, 6.3, 1.2, size=52, bold=True,
    color=MAROON, align='center')
txt(sl, 'Grade: ___  |  Adjective: _______________  |  n = 30',
    3.5, 3.15, 6.3, 0.55, size=14, color=DARK, align='center')
txt(sl, '(To be completed after Week 7 SUS administration)',
    3.5, 3.75, 6.3, 0.4, size=11, italic=True, color=DGREY, align='center')

txt(sl, 'SUS Score Interpretation Scale (Bangor et al., 2009):',
    0.5, 4.25, 12.3, 0.4, size=14, bold=True, color=BLUE)

interp = [
    ('≥ 90', 'A+', 'Best Imaginable', 'Acceptable', 'E8F5E9'),
    ('85–89', 'A', 'Excellent', 'Acceptable', 'E8F5E9'),
    ('80–84', 'B', 'Good', 'Acceptable', 'C8E6C9'),
    ('70–79', 'C', 'OK', 'Marginal', 'FFF9C4'),
    ('< 70', 'D/F', 'Poor – Awful', 'Not Acceptable', 'FFCCBC'),
]
col_iw = [1.5, 0.8, 2.2, 2.0, 0.05]
col_ix = [0.4, 1.95, 2.8, 5.05, 7.1]
for ri, (score, grade, adj, accept, color) in enumerate(interp):
    y = 4.75 + ri * 0.52
    highlight = (ri == 2)
    for ci, (val, cw, cx) in enumerate(zip([score, grade, adj, accept, '← TARGET' if highlight else ''],
                                            col_iw, col_ix)):
        card(sl, cx, y, cw, 0.46, fill=color if not highlight else 'A5D6A7',
             border_color='AAAAAA')
        txt(sl, val, cx + 0.05, y + 0.07, cw - 0.1, 0.32,
            size=12, bold=highlight, color=MAROON if (highlight and ci == 4) else DARK)

txt(sl, 'Respondents: 20 Parents (35–55)  +  10 Young Professionals (21–35)  |  '
    'Administered after guided live demo using SmartSpend Demo Mode',
    0.5, 7.1, 12.3, 0.32, size=11, italic=True, color=DGREY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — KEY FEATURES DEMO OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'Demo Overview — 8–9 Minute Flow', 0.4, 0.12, 12.5, 0.75,
    size=32, bold=True, color=WHITE)
slide_number(sl, 14)

demo_steps = [
    ('0:00', '30s', '🏠 Home Screen', 'FHS score card, spending summary, wallet card', 'EAD9E0'),
    ('0:30', '2min', '🤖 AI Chat', 'Voice → text → wallet update → PH advisory question', 'DAE8F5'),
    ('2:30', '1.5min', '📷 Smart Import', '4-mode sheet, Batch Screenshots demo (40+ platforms)', 'D9F0DA'),
    ('4:00', '1.5min', '💵 Wallets & Settings', 'Wallet card, Log Allowance, Balance Mode toggle', 'FFF8DC'),
    ('5:30', '1.5min', '📊 Analytics', 'Pie chart, 50/30/20, FHS trend, market insights', 'F3E5F5'),
    ('7:00', '1min', '🏥 FHS Breakdown', 'Tap score card, all 4 components, "Ask AI" link', 'FFE8D0'),
    ('8:00', '30s', '🎮 Gamification', 'Daily quests, mood check-in, achievement badges', 'E0F7FA'),
    ('8:30', '30s', '🎬 Architecture + Close', 'Serverless, context injection, Filipino-first design', 'F0E6FF'),
]
for i, (time, dur, title, desc, color) in enumerate(demo_steps):
    col = i % 2; row = i // 2
    x = 0.4 + col * 6.47
    y = 1.2 + row * 1.55
    card(sl, x, y, 6.23, 1.4, fill=color, border_color='CCCCCC')
    txt(sl, time, x + 0.1, y + 0.08, 0.85, 0.5, size=14, bold=True, color=MAROON)
    txt(sl, dur, x + 1.0, y + 0.08, 0.9, 0.4, size=11, color=DGREY, italic=True)
    txt(sl, title, x + 1.95, y + 0.08, 4.1, 0.5, size=14, bold=True, color=BLUE)
    txt(sl, desc, x + 0.1, y + 0.7, 6.0, 0.55, size=12, color=DARK)

card(sl, 0.4, 7.1, 12.5, 0.32, fill='5C0E24')
txt(sl, '📌  Before presenting: Reset AI limit (AI screen → ⋮ → Reset Daily Limit)  |  '
    'Load Demo Data if needed (Profile → Load Demo Data)',
    0.55, 7.13, 12.2, 0.26, size=11, bold=True, color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CONCLUSIONS & RECOMMENDATIONS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 1.0, MAROON)
accent_bar(sl, top=0.95, height=0.06)
txt(sl, 'Conclusions & Recommendations', 0.4, 0.12, 12.5, 0.75,
    size=32, bold=True, color=WHITE)
slide_number(sl, 15)

concl = [
    ('Objective 1\n(Assessment)', 'Confirmed: Filipinos face manual effort burden, '
     'irregular budgeting, and lack of proactive financial feedback — '
     'validating SmartSpend\'s core design rationale.', 'EAD9E0', MAROON),
    ('Objective 2\n(Development)', 'SmartSpend v2.9.10 fully developed: 31 agentic actions, '
     'dual-mode FHS, 5-provider failover, 40+ screenshot platforms, all free-tier stack. '
     'Gemini 3.1 Flash-Lite confirmed as optimal primary model.', 'DAE8F5', BLUE),
    ('Objective 3\n(SUS Evaluation)', 'To be completed after Week 7. '
     'Target: SUS ≥ 80 (Good, Acceptable per Bangor et al., 2009). '
     'Instrument administered to 30 respondents after guided demo.', 'D9F0DA', GREEN),
]
for i, (label, text, color, col) in enumerate(concl):
    y = 1.2 + i * 1.7
    card(sl, 0.4, y, 12.5, 1.55, fill=color, border_color='CCCCCC')
    card(sl, 0.4, y, 1.6, 1.55, fill=f'{col[0]:02X}{col[1]:02X}{col[2]:02X}')
    txt(sl, label, 0.4, y + 0.3, 1.6, 0.95, size=12, bold=True, color=WHITE, align='center')
    txt(sl, text, 2.1, y + 0.25, 10.6, 1.05, size=14, color=DARK)

txt(sl, 'Post-Capstone Recommendations:', 0.5, 6.35, 12.3, 0.38,
    size=14, bold=True, color=MAROON)
recs = ['Paluwagan tracker (highest priority — uniquely Filipino)',
        'Backend API proxy (Firebase Cloud Function for key security)',
        '15th/30th payday cycle awareness  |  Family wallet sharing',
        'Google Play Store submission  |  SQLite encryption (SQLCipher)']
for i, r in enumerate(recs):
    txt(sl, f'→  {r}', 0.6, 6.78 + i * 0.18, 12.0, 0.18,
        size=12, color=DARK)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CLOSING / THANK YOU
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, MAROON)
rect(sl, 0, 6.6, 13.33, 0.9, GOLD)

txt(sl, 'Thank You', 0.8, 1.5, 11.5, 1.8, size=72, bold=True, color=WHITE, align='center')
txt(sl, 'SmartSpend — v2.9.10  |  Lucid Frame', 0.8, 3.4, 11.5, 0.55,
    size=20, color=GOLD, align='center', bold=True)
txt(sl, 'Brix A. Directo  ·  Cyrille John M. Rubis  ·  Djaunathan Albert S. Madayag',
    0.8, 4.0, 11.5, 0.45, size=15, color=WHITE, align='center')
txt(sl, 'Lorma Colleges  —  CCSE, BSIT 4th Year  —  AY 2026–2027, 1st Semester',
    0.8, 4.5, 11.5, 0.4, size=13, color=LGREY, align='center')
txt(sl, 'Adviser: Janelli M. Mendez, DIT', 0.8, 4.95, 11.5, 0.35,
    size=13, color=GOLD, align='center')
txt(sl, 'github.com/Zushikina-kun/smartspend-app',
    0.8, 5.55, 11.5, 0.35, size=13, color=LGREY, align='center', italic=True)
txt(sl, 'Open for Questions', 0.8, 6.65, 11.5, 0.4,
    size=18, bold=True, color=MAROON, align='center')

# ── SAVE ───────────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f'  Defense slides saved: {OUT.name}  ({OUT.stat().st_size//1024} KB)')
print(f'  16 slides total:')
print(f'    1. Title  2. Problem  3. Objectives  4. What is SmartSpend')
print(f'    5. Tech Stack  6. Agentic AI  7. LLM Benchmarking  8. General vs Finance LLMs')
print(f'    9. FHS Formula  10. App Comparison  11. Smart Import  12. Methodology')
print(f'    13. SUS Results [placeholder]  14. Demo Flow  15. Conclusions  16. Thank You')
